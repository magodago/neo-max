# Documentos profesionales con imagenes (PPTX). Uso: SKILL:documentos_imagenes <tema>
DESCRIPTION = "Crea documentos o presentaciones con imagenes (PPTX). Uso: SKILL:documentos_imagenes <tema>"

def run(task: str = "", **kwargs) -> str:
    import re
    from pathlib import Path
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor
    except ImportError:
        return "Error: pip install python-pptx"
    tema = (task or "Documento").strip()[:200]
    slug = re.sub(r"[^\w\s-]", "", tema)[:40].strip().replace(" ", "_") or "doc"
    base = Path(__file__).resolve().parent.parent
    out_dir = base / "output" / "documentos_imagenes"
    out_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1))
    p = tx.text_frame.paragraphs[0]
    p.text = tema[:80]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0x1a, 0x1a, 0x2e)
    for i in range(1, 3):
        slide = prs.slides.add_slide(blank)
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.6))
        tx.text_frame.paragraphs[0].text = "Seccion " + str(i) + ": " + tema[:50]
        tx.text_frame.paragraphs[0].font.size = Pt(24)
    out_file = out_dir / (slug + ".pptx")
    prs.save(out_file)
    return "Documento creado: " + str(out_file) + ". Tema: " + tema + ". Abre el archivo o comparte."
